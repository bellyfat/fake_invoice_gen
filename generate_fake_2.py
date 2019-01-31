
### File to generate invoices from template ppt

# TODO
# set NUM_FILES
# set items


import pptx
import random
from pptx import Presentation
from random import randint
from random import seed, sample
from datetime import datetime
from faker import Faker
from pptx.util import Pt
from faker.utils.datetime_safe import date

# ---------------- CHANGE THIS ------------

PPT_FILE = "SrcInvoice_format_2.pptx"

START_FILE_INDEX = 1
END_FILE_INDEX = 2

MIN_NUM_ITEMS = 2
MAX_NUM_ITEMS =  3


items = ["Mango", "Banana", "Charger", "Bottle", "Phone", "Wire", "Board", "Light",
"toothbrush",
"computer",
"lace",
"book",
"ring",
"cookie jar",
"shirt",
"shampoo",
"sandal",
"toothbrush",
"video games",
"rusty nail",
"USB drive",
"plate",
"headphones",
"drawer",
]

# ---------------------------------------------------------------------


seed(datetime.now())
fake = Faker()

def get_text(i, j=None, font=None, size=None):
    frame = None
    if j is not None:
        frame = s.shapes[i].shapes[j].text_frame
    else:
        frame = s.shapes[i].text_frame

    if font:
        frame.paragraphs[0].font.name = font
    if size:
        frame.paragraphs[0].font.size = Pt(size)
    return frame

def get_rand_item(item_list):
    return item_list[randint(-1, len(item_list)-1)]

def get_rand_items(item_list, num):
    return sample(item_list, num)

def get_3num_code():
    return str(randint(100, 999)) + "-" + str(randint(100, 999)) + "-" + str(randint(100, 999))

def generate_adress():
    return get_rand_item(names) + '\n' +get_rand_item(cities) + '\n' + get_rand_item(states) + '\n' +  str(randint(100000, 999999))

def generate_company():
    return get_rand_item(company_names) +" Inc."+ '\n' + get_rand_item(cities) + '\n' + get_rand_item(states) + ' ' + str(randint(100000, 999999))


def generate_fake_phone_num():
    return str(randint(100, 999)) + "-" + str(randint(100, 999)) + "-" + str(randint(1000, 9999))

def generate_fake_company():
    return fake.company().split()[0]+" Inc."+ '\n' + fake.state() + ' ' + fake.postalcode()

def generate_fake_address():
    return fake.state() + ' ' + fake.postalcode()  + '\n' + fake.country_code(representation="alpha-3") + '\n\n' +  generate_fake_phone_num() + '\n' + fake.email()

def generate_vals():
    return "Invoice Number: "+str(randint(100000, 9999999))+"\nP.O/S.O: "+str(randint(100000, 9999999)) \
      + "\nInvoice Date: "+fake.date_this_decade(before_today=True, after_today=False).strftime("%b %d, %Y")+"\nPayment Date: "\
      + fake.date_this_decade(before_today=False, after_today=True).strftime("%b %d, %Y")+"\nAmount Due (USD): $"+str(format(randint(20, 500), '.2f'))



for i in range(START_FILE_INDEX, END_FILE_INDEX):

    prs = Presentation(PPT_FILE)
    s = prs.slides[0]
    
    # get the inv comp tf
    inv_comp_tf = get_text(4)
    # set the inv comp tf 
    inv_comp_tf_paragraphs = inv_comp_tf.paragraphs
    inv_comp_tf_paragraphs[0].text = fake.company().split()[0]+" Inc."
    inv_comp_tf_paragraphs[0].font.size = Pt(14)
    inv_comp_tf_paragraphs[0].font.name = 'Arial'
    inv_comp_tf_paragraphs[0].font.bold = True

    inv_comp_tf_paragraphs[1].text = fake.street_name()
    inv_comp_tf_paragraphs[1].font.size = Pt(14)
    inv_comp_tf_paragraphs[1].font.name = 'Arial'

    inv_comp_tf_paragraphs[2].text = fake.state() + ' ' + fake.postalcode()
    inv_comp_tf_paragraphs[2].font.size = Pt(14)
    inv_comp_tf_paragraphs[2].font.name = 'Arial'

    # get the bill to text field
    bill_to_tf = get_text(5)
    # set the bill to text field
    bill_to_tf_paragraphs = bill_to_tf.paragraphs
    bill_to_tf_paragraphs[1].text = fake.name().split()[0]
    bill_to_tf_paragraphs[1].font.size = Pt(14)
    bill_to_tf_paragraphs[1].font.name = 'Arial'
    bill_to_tf_paragraphs[1].font.bold = True

    bill_to_tf_paragraphs[2].text = fake.street_name()
    bill_to_tf_paragraphs[2].font.size = Pt(14)
    bill_to_tf_paragraphs[2].font.name = 'Arial'

    bill_to_tf_paragraphs[3].text = fake.state() + ' ' + fake.postalcode()
    bill_to_tf_paragraphs[3].font.size = Pt(14)
    bill_to_tf_paragraphs[3].font.name = 'Arial'

    bill_to_tf_paragraphs[4].text = fake.country_code(representation="alpha-3")
    bill_to_tf_paragraphs[4].font.size = Pt(14)
    bill_to_tf_paragraphs[4].font.name = 'Arial'

    bill_to_tf_paragraphs[6].text = generate_fake_phone_num()
    bill_to_tf_paragraphs[6].font.size = Pt(14)
    bill_to_tf_paragraphs[6].font.name = 'Arial'

    bill_to_tf_paragraphs[7].text = fake.email()
    bill_to_tf_paragraphs[7].font.size = Pt(14)
    bill_to_tf_paragraphs[7].font.name = 'Arial'

    # get tf right to invoice
    right_to_inv_tf = s.shapes[24].text_frame
    # set the value
    right_to_inv_tf_paragraphs = right_to_inv_tf.paragraphs
    right_to_inv_tf_paragraphs[0].text = str(i)
    right_to_inv_tf_paragraphs[0].font.size = Pt(14)
    right_to_inv_tf_paragraphs[0].font.name = 'Arial'

    month_dict = {1:'Jan', 2:'Feb',3:'Mar', 4:'Apr',5:'May',6:'Jun',7:'Jul',8:'Aug',9:'Sep',10:'Oct',11:'Nov',12:'Dec'}
    dt = fake.date(pattern="%m-%d-%Y", end_datetime="now")
    dtl = dt.split('-')
    # create date object date(yr, m, day)
    inv_date = date(int(dtl[2]), int(dtl[0]), int(dtl[1]))

    # get the payment date
    paymnt_date = fake.date_between_dates(date_start=inv_date, date_end=None) 
    # paymnt_date in (yr, m, day) fmt

    dtl[0] = month_dict[int(dtl[0])]
    paymnt_mnth = month_dict[int(paymnt_date.month)]

    inv_date_str = dtl[0] + ' ' + dtl[1] + ', ' + dtl[2]
    paymnt_date_str = str(paymnt_mnth) + ' ' + str(paymnt_date.day) + ', ' + str(paymnt_date.year)

    right_to_inv_tf_paragraphs[1].text = inv_date_str
    right_to_inv_tf_paragraphs[1].font.size = Pt(14)
    right_to_inv_tf_paragraphs[1].font.name = 'Arial'

    right_to_inv_tf_paragraphs[2].text = paymnt_date_str
    right_to_inv_tf_paragraphs[2].font.size = Pt(14)
    right_to_inv_tf_paragraphs[2].font.name = 'Arial'

    amt_due = round(random.uniform(0.1,1) * 100, 2)
    right_to_inv_tf_paragraphs[3].text = '$' + str(amt_due)
    right_to_inv_tf_paragraphs[3].font.size = Pt(14)
    right_to_inv_tf_paragraphs[3].font.name = 'Arial'

    # group shapes are at 10, 11, 15
    qty_gshapes = s.shapes[10].shapes
    name_gshapes = s.shapes[11].shapes
    amt_gshapes = s.shapes[15].shapes

    price_0_tf = s.shapes[13].text_frame
    price_1_tf = s.shapes[14].text_frame

    total_tf = s.shapes[18].text_frame
    amt_due_total_tf = s.shapes[19].text_frame
    sub_total_tf = s.shapes[27].text_frame

    # define the qty, price list and cal amt
    qty_list = [randint(1,5), randint(1,5)]
    price_list = [round(random.uniform(0.1,1)*100,2), round(random.uniform(0.1,1)*100,2)]
    amt_list = [0.0, 0.0]
    sub_total = 0.0
    for i in range(2):
        amt_list[i] = round(qty_list[i] * price_list[i],2)
        sub_total += amt_list[i]

    ca_9_percent = round(sub_total * 0.09, 2)
    total_val = round(sub_total + ca_9_percent, 2)
    total_amt_due = round(total_val + amt_due, 2)

    # now set the calculated value
    n = len(items)
    indx = randint(0, n)
    name_gshapes[0].text_frame.paragraphs[0].text = str(items[indx])
    name_gshapes[0].text_frame.paragraphs[0].font.size = Pt(14)
    name_gshapes[0].text_frame.paragraphs[0].font.name = 'Arial'
    name_gshapes[0].text_frame.paragraphs[0].font.bold = True

    name_gshapes[0].text_frame.paragraphs[1].text = 'test1'
    name_gshapes[0].text_frame.paragraphs[1].font.size = Pt(14)
    name_gshapes[0].text_frame.paragraphs[1].font.name = 'Arial'
    
    name_gshapes[1].text_frame.paragraphs[0].text = str(items[randint(0, n)])
    name_gshapes[1].text_frame.paragraphs[0].font.size = Pt(14)
    name_gshapes[1].text_frame.paragraphs[0].font.name = 'Arial'
    name_gshapes[1].text_frame.paragraphs[0].font.bold = True

    name_gshapes[1].text_frame.paragraphs[1].text = 'test2'
    name_gshapes[1].text_frame.paragraphs[1].font.size = Pt(14)
    name_gshapes[1].text_frame.paragraphs[1].font.name = 'Arial'    

    # set the quantity tf
    qty_gshapes[0].text = str(qty_list[0])
    qty_gshapes[0].text_frame.paragraphs[0].font.size = Pt(14)
    qty_gshapes[0].text_frame.paragraphs[0].font.name = 'Arial'

    qty_gshapes[1].text = str(qty_list[1])
    qty_gshapes[1].text_frame.paragraphs[0].font.size = Pt(14)
    qty_gshapes[1].text_frame.paragraphs[0].font.name = 'Arial'

    # set the price tf
    price_0_tf.text = '$' + str(price_list[0])
    price_0_tf.paragraphs[0].font.size = Pt(14)
    price_0_tf.paragraphs[0].font.name = 'Arial'
    
    price_1_tf.text = '$' + str(price_list[1])
    price_1_tf.paragraphs[0].font.size = Pt(14)
    price_0_tf.paragraphs[0].font.name = 'Arial'

    # set the amount gshapes
    amt_gshapes[0].text = '$' + str(amt_list[0])
    amt_gshapes[0].text_frame.paragraphs[0].font.size = Pt(14)
    amt_gshapes[0].text_frame.paragraphs[0].font.name = 'Arial'

    amt_gshapes[1].text = '$' + str(amt_list[1])
    amt_gshapes[1].text_frame.paragraphs[0].font.size = Pt(14)
    amt_gshapes[1].text_frame.paragraphs[0].font.name = 'Arial'

    # set sub total tf
    sub_total_tf.paragraphs[0].text = '$' + str(sub_total)
    sub_total_tf.paragraphs[0].font.size = Pt(14)
    sub_total_tf.paragraphs[0].font.name = 'Arial'

    sub_total_tf.paragraphs[1].text = '$' + str(ca_9_percent)
    sub_total_tf.paragraphs[1].font.size = Pt(14)
    sub_total_tf.paragraphs[1].font.name = 'Arial'

    # set total tf 
    total_tf.paragraphs[0].text = '$' + str(total_val)
    total_tf.paragraphs[0].font.size = Pt(14)
    total_tf.paragraphs[0].font.name = 'Arial'

    # set amt due tf
    amt_due_total_tf.paragraphs[0].text = '$' + str(total_amt_due)
    amt_due_total_tf.paragraphs[0].font.size = Pt(14)
    amt_due_total_tf.paragraphs[0].font.name = 'Arial'

    # print(bill_to_tf_paragraphs[1].text)

    # print(bill_to_tf.text)


    # pos = 0
    # for sh in s.shapes:
    #     # print(sh, pos)
    #     if(type(sh) is pptx.shapes.group.GroupShape):
    #         print(sh, pos)
    #     pos += 1
    

    # print('TEST_START')
    # print(len(s.shapes[5].text_frame.paragraphs))
    # print('TEST_END')

    # for x in s.shapes:
    #     print(s.shapes[pos], pos)
    #     if(x.has_text_frame):
    #         print('-----------------------------')
    #         print(x.text_frame.text)
    #         print(pos)
    #         print('-----------------------------')
    #     pos += 1


    # pending_payment = get_text(5)
    # bill_to_sec = get_text(9)
    # ship_to = get_text(10)
    # fedex = get_text(11)
    # terms_of_use = get_text(12)
    # table1 = s.shapes[13].table
    # table2 = s.shapes[14].table
    
    # print(bill_to.paragraphs[1].runs[1].text)
    
    # # test_code_to_fill_values
    # bill_to.paragraphs[1].text = generate_fake_company()
    # bill_to.paragraphs[1].font.size = Pt(14)
    # bill_to.paragraphs[1].font.name = 'Arial'
    ##### get frames ########
    # company = get_text(3)
    # bill_to = get_text(4)
    # headers = get_text(5)

    # total_vals = get_text(15)
    # amt_due_vals = get_text(16)

    ######## fill values ##############
    # company.text = generate_fake_company()
    # company.paragraphs[0].runs[0].font.bold=True
    # company.paragraphs[0].font.size=Pt(14)
    # company.paragraphs[0].font.name='Arial'
    # bill_to.text = generate_fake_adress()
    # bill_to.paragraphs[0].runs[0].font.bold=True
    # bill_to.paragraphs[0].font.size=Pt(14)
    # bill_to.paragraphs[0].font.name='Arial'
    # headers.text = generate_vals()
    # headers.paragraphs[0].font.size=Pt(14)
    # headers.paragraphs[0].font.name='Arial'

    # num_items = randint(MIN_NUM_ITEMS, MAX_NUM_ITEMS)
    # selected_items = get_rand_items(items, num_items)
    # sub_total = 0
    
    # print("---- "+str(i)+"-------")

    # for idx in range(num_items):
    #     qty = randint(1, 5)
    #     item = selected_items[idx]
    #     uprice = randint(1, 10)*5
    #     amt = qty*uprice
    #     sub_total += amt

    #     get_text(table_qty, idx, font='Arial', size=14).text = str(qty).zfill(2)
    #     get_text(table_name, idx, font='Arial', size=14).text = str(item)
    #     get_text(table_uprice, idx, font='Arial', size=14).text = "$"+str(format(uprice, '.2f'))
    #     get_text(table_amt, idx, font='Arial', size=14).text = "$"+str(format(amt, '.2f'))
    
    # total = sub_total

    # # subtotal_vals.text = str(format(sub_total,'.2f'))
    # # tax_vals.text = str( format(tax,'.2f'))
    # total_vals.text = "$"+str(format(total, '.2f'))
    # amt_due_vals.text = "$"+str(format(total+randint(0, 100), '.2f'))

    # amt_due_vals.paragraphs[0].font.bold = True
    # amt_due_vals.paragraphs[0].font.name = 'Arial Narrow'
    # total_vals.paragraphs[0].font.name = 'Arial Narrow'

    prs.save("fmt2//T2_RC_"+str(i)+".pptx")

